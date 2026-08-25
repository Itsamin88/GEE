"""Word, PowerPoint, OpenDocument and legacy Office extraction.

Paragraphs, headings, tables, captions, hyperlinks, embedded images and
document metadata are all kept, with the document-level provenance intact
(brief §20).
"""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass, field
from typing import Any

from ..logging_setup import get_logger

log = get_logger("office")


@dataclass
class OfficeImage:
    data: bytes
    name: str
    extension: str
    width: int = 0
    height: int = 0
    nearby_text: str = ""


@dataclass
class OfficeResult:
    ok: bool = False
    text: str = ""
    headings: list[str] = field(default_factory=list)
    paragraphs: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)
    hyperlinks: list[str] = field(default_factory=list)
    images: list[OfficeImage] = field(default_factory=list)
    metadata: dict[str, str] = field(default_factory=dict)
    parser: str = ""
    parser_status: str = "not_attempted"
    text_status: str = "not_attempted"
    table_status: str = "not_attempted"
    image_status: str = "not_attempted"
    detail: str = ""


def extract_docx(data: bytes, *, extract_images: bool = True,
                 min_image_pixels: int = 90000) -> OfficeResult:
    result = OfficeResult(parser="python-docx")
    try:
        import docx  # type: ignore
    except ImportError:
        result.parser_status = "unsupported_format"
        result.detail = "python-docx is not installed"
        return result
    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as exc:
        result.parser_status = "corrupt"
        result.text_status = "failed"
        result.detail = f"{type(exc).__name__}: {exc}"
        return result

    try:
        props = document.core_properties
        for key in ("title", "author", "subject", "created", "modified", "category",
                    "comments", "keywords", "last_modified_by"):
            value = getattr(props, key, None)
            if value:
                result.metadata[key] = str(value)[:500]
    except Exception:
        pass

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        result.paragraphs.append(text)
        style = (paragraph.style.name if paragraph.style is not None else "") or ""
        if style.lower().startswith("heading") or style.lower() in ("title", "subtitle"):
            result.headings.append(text[:300])

    for table in document.tables:
        rows: list[list[str]] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(cells)
        if len(rows) > 1:
            result.tables.append(rows[:400])

    try:
        for rel in document.part.rels.values():
            if rel.reltype.endswith("/hyperlink") and getattr(rel, "is_external", False):
                result.hyperlinks.append(str(rel.target_ref))
    except Exception:
        pass

    if extract_images:
        result.images = _zip_images(data, min_image_pixels)
        result.image_status = "extracted" if result.images else "none_found"
    else:
        result.image_status = "not_attempted"

    result.text = "\n".join(result.paragraphs)
    for rows in result.tables:
        result.text += "\n\n" + "\n".join(" | ".join(r) for r in rows)
    result.parser_status = "parsed"
    result.text_status = "extracted" if result.text.strip() else "empty"
    result.table_status = "extracted" if result.tables else "none_found"
    result.ok = True
    return result


def extract_pptx(data: bytes, *, extract_images: bool = True,
                 min_image_pixels: int = 90000) -> OfficeResult:
    """Slide decks: text frames, notes and tables, slide by slide."""
    result = OfficeResult(parser="python-pptx")
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        # Fall back to reading the OOXML parts directly rather than failing.
        return _extract_ooxml_text(data, result, part_prefix="ppt/slides/",
                                   extract_images=extract_images,
                                   min_image_pixels=min_image_pixels)
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        result.parser_status = "corrupt"
        result.detail = f"{type(exc).__name__}: {exc}"
        return result

    for index, slide in enumerate(presentation.slides, start=1):
        pieces: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    pieces.append(text)
            if getattr(shape, "has_table", False):
                rows = [[c.text.strip() for c in row.cells] for row in shape.table.rows]
                rows = [r for r in rows if any(r)]
                if len(rows) > 1:
                    result.tables.append(rows)
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    pieces.append(f"[notes] {notes}")
        except Exception:
            pass
        if pieces:
            result.paragraphs.append(f"[slide {index}]\n" + "\n".join(pieces))
            result.headings.append(pieces[0][:200])

    if extract_images:
        result.images = _zip_images(data, min_image_pixels)
        result.image_status = "extracted" if result.images else "none_found"
    result.text = "\n\n".join(result.paragraphs)
    result.parser_status = "parsed"
    result.text_status = "extracted" if result.text.strip() else "empty"
    result.table_status = "extracted" if result.tables else "none_found"
    result.ok = True
    return result


def extract_odf(data: bytes, *, extract_images: bool = True,
                min_image_pixels: int = 90000) -> OfficeResult:
    """OpenDocument text/presentation via its XML, without an extra dependency."""
    result = OfficeResult(parser="odf-xml")
    try:
        from bs4 import BeautifulSoup

        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
            if "content.xml" not in names:
                result.parser_status = "corrupt"
                result.detail = "no content.xml in the ODF container"
                return result
            content = archive.read("content.xml").decode("utf-8", "replace")
            if "meta.xml" in names:
                meta_soup = BeautifulSoup(archive.read("meta.xml").decode("utf-8", "replace"), "xml")
                for tag in meta_soup.find_all(True):
                    if tag.name in ("title", "creator", "date", "subject", "keyword"):
                        result.metadata[tag.name] = tag.get_text(strip=True)[:500]
        soup = BeautifulSoup(content, "xml")
        for node in soup.find_all(re.compile(r"(^|:)h$")):
            text = node.get_text(" ", strip=True)
            if text:
                result.headings.append(text[:300])
        for node in soup.find_all(re.compile(r"(^|:)p$")):
            text = node.get_text(" ", strip=True)
            if text:
                result.paragraphs.append(text)
        for table in soup.find_all(re.compile(r"(^|:)table$")):
            rows = []
            for row in table.find_all(re.compile(r"(^|:)table-row$")):
                cells = [c.get_text(" ", strip=True)
                         for c in row.find_all(re.compile(r"(^|:)table-cell$"))]
                if any(cells):
                    rows.append(cells)
            if len(rows) > 1:
                result.tables.append(rows[:400])
    except zipfile.BadZipFile as exc:
        result.parser_status = "corrupt"
        result.detail = f"bad ODF container: {exc}"
        return result
    except Exception as exc:
        result.parser_status = "corrupt"
        result.detail = f"{type(exc).__name__}: {exc}"
        return result

    if extract_images:
        result.images = _zip_images(data, min_image_pixels, prefix="Pictures/")
        result.image_status = "extracted" if result.images else "none_found"
    result.text = "\n".join(result.paragraphs)
    result.parser_status = "parsed"
    result.text_status = "extracted" if result.text.strip() else "empty"
    result.table_status = "extracted" if result.tables else "none_found"
    result.ok = True
    return result


def extract_legacy_doc(data: bytes) -> OfficeResult:
    """Legacy binary .doc — best-effort text recovery, honestly labelled."""
    result = OfficeResult(parser="legacy-doc")
    try:
        import olefile  # type: ignore
    except ImportError:
        result.parser_status = "unsupported_format"
        result.detail = "olefile is not installed; the original file is preserved unparsed"
        return result
    try:
        if not olefile.isOleFile(io.BytesIO(data)):
            result.parser_status = "corrupt"
            result.detail = "not an OLE container"
            return result
        ole = olefile.OleFileIO(io.BytesIO(data))
        stream_names = ["WordDocument", "1Table", "0Table"]
        chunks: list[str] = []
        for name in stream_names:
            if ole.exists(name):
                raw = ole.openstream(name).read()
                # Recover printable runs; a binary .doc has no public simple format.
                for match in re.finditer(rb"(?:[\x20-\x7e\xc0-\xff]{6,})", raw):
                    chunks.append(match.group().decode("latin-1", "ignore"))
        text = "\n".join(chunks)
        result.text = re.sub(r"\n{3,}", "\n\n", text)[:400000]
        result.parser_status = "parsed" if result.text.strip() else "corrupt"
        result.text_status = "extracted" if result.text.strip() else "empty"
        result.detail = "best-effort recovery from a legacy binary .doc; verify against the original"
        result.ok = bool(result.text.strip())
    except Exception as exc:
        result.parser_status = "corrupt"
        result.detail = f"{type(exc).__name__}: {exc}"
    return result


def _extract_ooxml_text(data: bytes, result: OfficeResult, *, part_prefix: str,
                        extract_images: bool, min_image_pixels: int) -> OfficeResult:
    try:
        from bs4 import BeautifulSoup

        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            parts = sorted(n for n in archive.namelist() if n.startswith(part_prefix)
                           and n.endswith(".xml"))
            for name in parts:
                soup = BeautifulSoup(archive.read(name).decode("utf-8", "replace"), "xml")
                text = " ".join(t.get_text(" ", strip=True) for t in soup.find_all(re.compile(r"(^|:)t$")))
                if text.strip():
                    result.paragraphs.append(text.strip())
    except Exception as exc:
        result.parser_status = "corrupt"
        result.detail = f"{type(exc).__name__}: {exc}"
        return result
    if extract_images:
        result.images = _zip_images(data, min_image_pixels)
        result.image_status = "extracted" if result.images else "none_found"
    result.text = "\n\n".join(result.paragraphs)
    result.parser = "ooxml-xml"
    result.parser_status = "parsed"
    result.text_status = "extracted" if result.text.strip() else "empty"
    result.ok = True
    return result


def _zip_images(data: bytes, min_pixels: int, prefix: str = "") -> list[OfficeImage]:
    """Pull embedded media out of any OOXML/ODF container, safely."""
    images: list[OfficeImage] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            for info in archive.infolist():
                name = info.filename
                if prefix and not name.startswith(prefix):
                    if "/media/" not in name:
                        continue
                elif "/media/" not in name and not name.startswith(prefix):
                    continue
                extension = name.rsplit(".", 1)[-1].lower() if "." in name else ""
                if extension not in ("jpg", "jpeg", "png", "gif", "tif", "tiff", "webp", "bmp", "emf", "wmf"):
                    continue
                if info.file_size > 30_000_000:
                    continue
                blob = archive.read(name)
                width, height = _dimensions(blob)
                if width * height and width * height < min_pixels:
                    continue
                images.append(
                    OfficeImage(data=blob, name=name.rsplit("/", 1)[-1],
                                extension=extension, width=width, height=height)
                )
    except Exception as exc:
        log.debug("embedded media extraction failed: %s", exc)
    return images


def _dimensions(blob: bytes) -> tuple[int, int]:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(blob)) as img:
            return img.width, img.height
    except Exception:
        return 0, 0
